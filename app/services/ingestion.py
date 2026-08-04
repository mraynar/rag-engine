from typing import Optional
import re
import time
import random
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from threading import Lock

import chromadb
from google.genai import types

from app.core.config import (
    VECTOR_STORE_DIR,
    get_embedding_model,
    get_gemini_client,
    get_generation_model,
)


CHUNK_SIZE = 300


def _split_sentences(text: str) -> list[str]:
    text = text.strip()
    raw = re.split(r"(?<=[.!?])\s+|\n\s*\n", text)
    return [s.strip() for s in raw if s.strip()]


def chunk_by_sentences(text: str, chunk_size: int = CHUNK_SIZE) -> list[str]:
    sentences = _split_sentences(text)
    chunks: list[str] = []
    current = ""

    for sentence in sentences:
        # oversized sentence becomes its own chunk
        if len(sentence) > chunk_size:
            if current:
                chunks.append(current.strip())
                current = ""
            chunks.append(sentence)
            continue

        candidate = f"{current} {sentence}".strip() if current else sentence
        if len(candidate) > chunk_size and current:
            chunks.append(current.strip())
            current = sentence
        else:
            current = candidate

    if current:
        chunks.append(current.strip())

    return chunks


def parse_txt(file_path: Path) -> list[str]:
    text = _read_text_with_fallback_encoding(file_path)
    return chunk_by_sentences(text)


def _read_text_with_fallback_encoding(file_path: Path) -> str:
    for enc in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return file_path.read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    raise ValueError(f"Gagal membaca file '{file_path.name}' — encoding tidak dikenali")


def parse_docx(file_path: Path) -> list[str]:
    from docx import Document  # python-docx
    doc = Document(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    combined = "\n".join(paragraphs)
    return chunk_by_sentences(combined)


def parse_pdf(file_path: Path) -> list[str]:
    import pdfplumber
    pages_text: list[str] = []
    with pdfplumber.open(str(file_path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text and text.strip():
                pages_text.append(text.strip())
    combined = "\n\n".join(pages_text)
    return chunk_by_sentences(combined)


def parse_xlsx(file_path: Path) -> list[str]:
    import pandas as pd
    xls = pd.read_excel(str(file_path), sheet_name=None)
    chunks = []
    for sheet_name, df in xls.items():
        sheet_chunks = _dataframe_to_chunks(df)
        context_chunks = [f"[Sheet: {sheet_name}] {chunk}" for chunk in sheet_chunks]
        chunks.extend(context_chunks)
    return chunks


def parse_csv(file_path: Path) -> list[str]:
    import pandas as pd
    df = None
    last_error = None
    for enc in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            df = pd.read_csv(str(file_path), encoding=enc)
            break
        except UnicodeDecodeError as e:
            last_error = e
            continue
    if df is None:
        raise ValueError(f"Gagal membaca file CSV — encoding tidak dikenali: {last_error}")
    return _dataframe_to_chunks(df)


def _dataframe_to_chunks(df) -> list[str]:
    """one row = one chunk; row-based because each table row is a standalone fact."""
    chunks: list[str] = []
    for _, row in df.iterrows():
        parts = []
        for col, val in row.items():
            if val is None:
                continue
            str_val = str(val).strip()
            if str_val and str_val.lower() not in ("nan", "none", ""):
                parts.append(f"{col}: {str_val}")
        if parts:
            chunks.append(", ".join(parts))
    return chunks


def parse_pptx(file_path: Path) -> list[str]:
    from pptx import Presentation  # python-pptx
    prs = Presentation(str(file_path))
    chunks: list[str] = []
    for slide in prs.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            chunks.append("\n".join(texts))
    return chunks


def parse_image(file_path: Path) -> list[str]:
    extension = file_path.suffix.lower()
    mime_type = "image/png"
    if extension in (".jpg", ".jpeg"):
        mime_type = "image/jpeg"
    elif extension == ".webp":
        mime_type = "image/webp"

    content = file_path.read_bytes()
    part = types.Part.from_bytes(data=content, mime_type=mime_type)

    prompt = (
        "Lakukan analisis mendalam terhadap gambar ini.\n"
        "1. Ekstrak dan tuliskan semua teks (OCR) yang terlihat pada gambar dengan sangat akurat dan lengkap.\n"
        "2. Jelaskan elemen visual, tabel, diagram, atau grafik yang ada di dalam gambar beserta nilainya secara detail.\n"
        "3. Berikan deskripsi ringkas tentang keseluruhan isi atau konteks gambar.\n"
        "Format output berupa teks markdown terstruktur agar mudah dibaca."
    )

    try:
        response = get_gemini_client().models.generate_content(
            model=get_generation_model(),
            contents=[part, prompt],
        )
        text = response.text or ""
    except Exception as e:
        raise RuntimeError(f"Gagal menganalisis gambar menggunakan model Gemini: {e}")

    return chunk_by_sentences(text)


_PARSERS = {
    ".txt":  parse_txt,
    ".docx": parse_docx,
    ".pdf":  parse_pdf,
    ".xlsx": parse_xlsx,
    ".xls":  parse_xlsx,   # pandas read_excel handles .xls too
    ".csv":  parse_csv,
    ".pptx": parse_pptx,
    ".jpg":  parse_image,
    ".jpeg": parse_image,
    ".png":  parse_image,
    ".webp": parse_image,
}

_SUPPORTED_FORMATS = ", ".join(sorted(_PARSERS.keys()))

_request_times: list[float] = []
_request_lock = Lock()
MAX_RPM = 12  # strict limit for Gemini Free Tier (15 RPM)


def _wait_for_rate_limit_slot():
    with _request_lock:
        now = time.time()
        while _request_times and _request_times[0] < now - 60:
            _request_times.pop(0)
        if len(_request_times) >= MAX_RPM:
            sleep_time = 60 - (now - _request_times[0]) + 0.1
            time.sleep(max(sleep_time, 0.0))
        _request_times.append(time.time())


def _embed_texts_batch(
    texts: list[str],
    batch_size: int = 50,
    max_workers: int = 1,
) -> list[list[float]]:
    """Embed texts in parallel batches; result order matches input order."""
    batches = [texts[i : i + batch_size] for i in range(0, len(texts), batch_size)]
    total_batches = len(batches)
    print(f"[ingestion] Starting embedding: {len(texts)} chunks → {total_batches} batches of max {batch_size}")

    # Exponential backoff schedule (seconds base, excluding jitter 0-5s):
    # attempt 0: ~10s, 1: ~30s, 2: ~70s, 3: ~130s, 4: ~250s, 5: ~490s
    _BACKOFF_SCHEDULE = [10, 30, 70, 130, 250, 490]
    _MAX_RETRIES = len(_BACKOFF_SCHEDULE)

    def embed_one_batch(batch_idx_and_batch: tuple) -> list[list[float]]:
        batch_idx, batch = batch_idx_and_batch
        _wait_for_rate_limit_slot()
        time.sleep(1.0)  # fixed pacing on top of the rate-limit window
        for attempt in range(_MAX_RETRIES):
            try:
                result = get_gemini_client().models.embed_content(
                    model=get_embedding_model(),
                    contents=batch,
                    config=types.EmbedContentConfig(task_type="RETRIEVAL_DOCUMENT"),
                )
                print(f"[ingestion] Batch {batch_idx + 1}/{total_batches} OK ({len(batch)} chunks)")
                return [e.values for e in result.embeddings]
            except Exception as e:
                if "RESOURCE_EXHAUSTED" in str(e) and attempt < _MAX_RETRIES - 1:
                    backoff = _BACKOFF_SCHEDULE[attempt]
                    sleep_time = backoff + random.uniform(0, 5)
                    print(
                        f"[ingestion] Gemini Rate Limit (RESOURCE_EXHAUSTED) hit on batch "
                        f"{batch_idx + 1}/{total_batches}, attempt {attempt + 1}/{_MAX_RETRIES}. "
                        f"Sleeping for {sleep_time:.2f}s before retry..."
                    )
                    time.sleep(sleep_time)
                    continue
                raise
        return []  # unreachable

    all_embeddings: list[list[list[float]] | None] = [None] * total_batches
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_idx = {
            executor.submit(embed_one_batch, (i, batch)): i
            for i, batch in enumerate(batches)
        }
        for future in as_completed(future_to_idx):
            idx = future_to_idx[future]
            all_embeddings[idx] = future.result()  # re-raises exception from thread

    flat: list[list[float]] = []
    for batch_result in all_embeddings:
        flat.extend(batch_result)  # type: ignore[arg-type]
    print(f"[ingestion] Embedding complete: {len(flat)} embeddings total.")
    return flat


def ingest_document(file_path: Path, filename: str, category: Optional[str] = None) -> int:
    """Parse, embed, and index a document into ChromaDB. Returns chunk count."""
    extension = file_path.suffix.lower()
    parser = _PARSERS.get(extension)
    if parser is None:
        raise ValueError(
            f"Format file '{extension}' belum didukung. "
            f"Format yang didukung: {_SUPPORTED_FORMATS}"
        )

    chunks = parser(file_path)
    if not chunks:
        return 0

    chroma_client = chromadb.PersistentClient(path=str(VECTOR_STORE_DIR))
    collection = chroma_client.get_or_create_collection(name="tps_docs")

    # delete stale chunks for this filename before re-indexing
    try:
        collection.delete(where={"source": filename})
    except Exception:
        pass

    embeddings = _embed_texts_batch(chunks)

    chunk_id_prefix = filename.replace(".", "_").replace(" ", "_")
    ids       = [f"{chunk_id_prefix}_{i}" for i in range(len(chunks))]
    metadatas = []
    for _ in chunks:
        meta = {"source": filename}
        if category is not None:
            meta["category"] = category
        metadatas.append(meta)

    collection.add(
        ids=ids,
        embeddings=embeddings,
        documents=chunks,
        metadatas=metadatas,
    )

    return len(chunks)


def delete_category_vector_data(category_name: str) -> None:
    try:
        chroma_client = chromadb.PersistentClient(path=str(VECTOR_STORE_DIR))
        collection = chroma_client.get_or_create_collection(name="tps_docs")
        collection.delete(where={"category": category_name})
    except Exception as e:
        print(f"[ingestion] Warning: failed to delete vector data for category '{category_name}': {e}")

